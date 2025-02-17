import os
import hashlib
import datetime
import magic
import re
import sqlite3
from pathlib import Path
import threading
from queue import Queue
import concurrent.futures
from typing import Dict, List, Set, Optional, Tuple
import sys
import time
import logging
from dataclasses import dataclass
from contextlib import contextmanager
import json

import tempfile
import pdfminer.high_level
import easyocr
import csv
import pandas as pd
from pptx import Presentation
from pdf2image import convert_from_path
from docx import Document
import magic
import win32com.client

class TextExtractor:
    def __init__(self, file_path: str):
        self.file_path = file_path
        self.reader_ocr = None  # Initialize on demand to avoid pickling issues

    def _get_ocr_reader(self):
        if self.reader_ocr is None:
            self.reader_ocr = easyocr.Reader(['en'], gpu=False)
        return self.reader_ocr

    def extract_text(self) -> str:
        ext = os.path.splitext(self.file_path)[1].lower()
        try:
            if ext == '.pdf':
                return self.extract_pdf()
            elif ext == '.docx':
                return self.extract_docx()
            elif ext in ['.xls', '.xlsx']:
                return self.extract_excel()
            elif ext in ['.csv']:
                return self.extract_csv()
            elif ext in ['.pptx', '.ppt']:
                return self.extract_text_from_ppt()
            elif ext in [".jpg", ".jpeg", ".png", ".bmp", ".gif"]:
                return self.extract_image()
            elif ext in [".txt", ".md", ".html", ".htm", ".xml", ".json"]:
                return self.extract_plain_text()
            else:
                return ""
        except Exception as e:
            logging.error(f"Error processing {self.file_path}: {e}")
            return ""

    

    def extract_pdf(self) -> str:
        try:
            # Try extracting text with pdfminer first
            text = pdfminer.high_level.extract_text(self.file_path)
            if text.strip():
                return text

            # Fallback: Convert PDF pages to images for OCR
            images = convert_from_path(self.file_path, dpi=300)
            ocr_text = []
            reader = self._get_ocr_reader()
            for i, image in enumerate(images):
                temp_img_path = f"temp_page_{i}.png"
                image.save(temp_img_path, "PNG")
                ocr_result = reader.readtext(temp_img_path, detail=0)
                if ocr_result:
                    ocr_text.append("\n".join(ocr_result))
                os.remove(temp_img_path)
            return "\n".join(ocr_text)
        except Exception as e:
            logging.error(f"PDF extraction failed: {e}")
            return ""



    def extract_docx(self) -> str:
        try:
            doc = Document(self.file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            logging.error(f"DOCX extraction failed: {e}")
            return ""

    def extract_excel(self) -> str:
        try:
            xl = pd.ExcelFile(self.file_path)
            return "\n".join(" ".join(map(str, row)) for sheet in xl.sheet_names for row in xl.parse(sheet).values)
        except Exception as e:
            logging.error(f"Excel extraction failed: {e}")
            return ""

    def extract_csv(self) -> str:
        try:
            with open(self.file_path, newline='', encoding='utf-8') as csvfile:
                reader = csv.reader(csvfile)
                return "\n".join(" ".join(row) for row in reader)
        except Exception as e:
            logging.error(f"CSV extraction failed: {e}")
            return ""



    def extract_text_from_ppt(self): 
        """Extract text from PowerPoint files including images.""" 
        import os
        import tempfile
        import logging
        import pythoncom 

        ext = os.path.splitext(self.file_path)[1].lower() 
        text = "" 
        
        if ext == ".pptx":         
            prs = Presentation(self.file_path) 
            
            # Extract text from shapes 
            for slide in prs.slides: 
                for shape in slide.shapes: 
                    if hasattr(shape, "text"): 
                        text += shape.text + "\n" 
            
            # Try to extract images and perform OCR 
            try: 
                reader = self._get_ocr_reader()  # Assuming self.reader_ocr is an EasyOCR reader instance
                
                with tempfile.TemporaryDirectory() as temp_dir: 
                    img_count = 0 
                    for slide in prs.slides: 
                        for shape in slide.shapes: 
                            if shape.shape_type == 13:  # Picture shape type 
                                image = shape.image 
                                img_path = os.path.join(temp_dir, f"img_{img_count}.png") 
                            
                                with open(img_path, "wb") as f: 
                                    f.write(image.blob) 
                                
                                # Perform OCR 
                                ocr_result = reader.readtext(img_path, detail=0) 
                                if ocr_result: 
                                    text += "\n" + "\n".join(ocr_result) 
                                img_count += 1 
                                
            except ImportError: 
                logging.warning("EasyOCR not installed. Image text extraction will be skipped.") 
            except Exception as e: 
                logging.error(f"Error extracting images from PPTX: {e}") 
                    
        elif ext == ".ppt": 
            
            try: 
                # Initialize COM before creating PowerPoint instance
                pythoncom.CoInitialize()  # Add this line
                
                # Dispatch PowerPoint application
                powerpoint = win32com.client.Dispatch("PowerPoint.Application") 
                powerpoint.Visible = True  
                
                presentation = powerpoint.Presentations.Open(os.path.abspath(self.file_path), WithWindow=False) 
                
                # Extract text from slides 
                for slide in presentation.Slides: 
                    for shape in slide.Shapes: 
                        if shape.HasTextFrame: 
                            if shape.TextFrame.HasText: 
                                text += shape.TextFrame.TextRange.Text + "\n" 
                
                # Save as temporary PPTX to extract images 
                with tempfile.TemporaryDirectory() as temp_dir: 
                    temp_pptx = os.path.join(temp_dir, "temp.pptx") 
                    presentation.SaveAs(temp_pptx) 
                    presentation.Close() 
                    powerpoint.Quit() 
                    
                    # Now extract images from the PPTX version 
                    original_file_path = self.file_path  # store original file path 
                    self.file_path = temp_pptx 
                    text += self.extract_text_from_ppt() 
                    self.file_path = original_file_path  # restore original path
                    
            except Exception as e: 
                logging.error(f"Error processing PPT file: {e}") 
                if 'presentation' in locals(): 
                    presentation.Close() 
                if 'powerpoint' in locals(): 
                    powerpoint.Quit()
            finally:
                # Uninitialize COM when done
                pythoncom.CoUninitialize()  # Add this line
        else: 
            raise ValueError(f"Unsupported PowerPoint file type: {ext}") 
            
        return text


    def extract_image(self) -> str:
        try:
            import cv2
            import numpy as np
            
            # First try to load the image with OpenCV
            try:
                img = cv2.imread(self.file_path)
                if img is None:
                    logging.error(f"Unable to load image with OpenCV: {self.file_path}")
                    raise ValueError(f"OpenCV could not load image: {self.file_path}")
                    
                # Convert from BGR to RGB (EasyOCR expects RGB)
                img_rgb = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
                
                # Pass the loaded image to EasyOCR instead of the file path
                reader = self._get_ocr_reader()
                return "\n".join(reader.readtext(img_rgb, detail=0))
                
            except Exception as cv_error:
                logging.error(f"OpenCV loading failed: {cv_error}")
                
                # Try PIL as fallback
                try:
                    from PIL import Image
                    image = Image.open(self.file_path)
                    img_array = np.array(image)
                    
                    reader = self._get_ocr_reader()
                    return "\n".join(reader.readtext(img_array, detail=0))
                    
                except Exception as pil_error:
                    logging.error(f"PIL loading failed: {pil_error}")
                    raise ValueError(f"All image loading methods failed for {self.file_path}")
                    
        except Exception as e:
            logging.error(f"Image extraction failed: {e}")
            return ""

    def extract_plain_text(self) -> str:
        try:
            with open(self.file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logging.error(f"Plain text extraction failed: {e}")
            return ""


@dataclass
class FileInfo:
    """Data class to store file information"""
    path: str
    is_directory: bool
    parent_directory: str
    size: int
    created_time: datetime.datetime
    modified_time: datetime.datetime
    accessed_time: datetime.datetime
    file_hash: Optional[str] = None
    mime_type: Optional[str] = None
    has_sensitive_data: bool = False
    patterns: str = ''
    sensitive_data: str = ''


class FileSystemIndexer:
    PATTERNS = {}
    with open('patterns.json', 'r') as file:
        data = json.load(file) 
    PATTERNS = data
    CHUNK_SIZE = 8192

    def __init__(self, db_path: str, log_level: int = logging.INFO):
        self.db_path = db_path
        self.setup_logging(log_level)
        self.setup_database()
        self.lock = threading.Lock()
        self.stats = {
            'processed_files': 0,
            'total_files': 0,
            'sensitive_files': 0,
            'total_size': 0
        }
        # Initialize thread local storage properly
        self._local = threading.local()
        self._local.conn = None
        self.extracted_text = TextExtractor("")

    def setup_logging(self, log_level: int) -> None:
        logging.basicConfig(
            level=log_level,
            format='%(asctime)s - %(levelname)s - %(message)s',
            datefmt='%Y-%m-%d %H:%M:%S'
        )
        self.logger = logging.getLogger(__name__)

    @contextmanager
    def get_db_connection(self):
        """Thread-safe database connection context manager with proper initialization"""
        if getattr(self._local, 'conn', None) is None:
            self._local.conn = sqlite3.connect(self.db_path)
            self._local.conn.execute('PRAGMA journal_mode=WAL')
            self._local.conn.execute('PRAGMA synchronous=NORMAL')
        
        try:
            yield self._local.conn
        except Exception as e:
            self.logger.error(f"Database error: {e}")
            if self._local.conn:
                self._local.conn.rollback()
            raise
        else:
            if self._local.conn:
                self._local.conn.commit()

    def setup_database(self) -> None:
        """Initialize SQLite database with optimized schema"""
        with sqlite3.connect(self.db_path) as conn:  # Use direct connection for setup
            conn.executescript('''
                CREATE TABLE IF NOT EXISTS files (
                    id INTEGER PRIMARY KEY,
                    path TEXT UNIQUE,
                    file_hash TEXT,
                    mime_type TEXT,
                    size INTEGER,
                    created_time TIMESTAMP,
                    modified_time TIMESTAMP,
                    accessed_time TIMESTAMP,
                    is_directory BOOLEAN,
                    parent_directory TEXT,
                    has_sensitive_data BOOLEAN,
                    patterns TEXT,
                    sensitive_data TEXT
                );
                
                CREATE INDEX IF NOT EXISTS idx_path ON files(path);
                CREATE INDEX IF NOT EXISTS idx_parent ON files(parent_directory);
                CREATE INDEX IF NOT EXISTS idx_sensitive ON files(has_sensitive_data);
                
                CREATE TABLE IF NOT EXISTS sensitive_data (
                    id INTEGER PRIMARY KEY,
                    file_id INTEGER,
                    file_path TEXT,
                    data_type TEXT,
                    line_number INTEGER,
                    matched_content TEXT,
                    FOREIGN KEY (file_id) REFERENCES files(id)
                );
                
                CREATE INDEX IF NOT EXISTS idx_file_id ON sensitive_data(file_id);
            ''')

    def calculate_file_hash(self, file_path: str) -> Optional[str]:
        """Calculate SHA-256 hash of file with improved performance"""
        sha256_hash = hashlib.sha256()
        try:
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(self.CHUNK_SIZE), b""):
                    sha256_hash.update(chunk)
            return sha256_hash.hexdigest()
        except (IOError, OSError) as e:
            self.logger.warning(f"Failed to calculate hash for {file_path}: {e}")
            return None

    def detect_sensitive_data(self, content: str) -> Tuple[str, str, List[Dict]]:
        findings = []
        patterns_found = set()
        line_numbers = set()

        for line_num, line in enumerate(content.split('\n'), 1):
            for pattern_name, pattern in self.PATTERNS.items():
                for match in re.finditer(pattern, line):
                    patterns_found.add(pattern_name)
                    line_numbers.add(line_num)
                    matched_content = {k: v for k, v in match.groupdict().items() if v}
                    findings.append({
                        'pattern': pattern_name,
                        'line_number': line_num,
                        'matched_content': str(matched_content)
                    })

        return (
            ','.join(map(str, sorted(line_numbers))),
            ','.join(sorted(patterns_found)),
            findings
        )
    
    def process_entry(self, entry: os.DirEntry) -> Optional[FileInfo]:
        try:
            stat = entry.stat(follow_symlinks=False)
            is_directory = entry.is_dir(follow_symlinks=False)
            
            file_info = FileInfo(
                path=entry.path,
                is_directory=is_directory,
                parent_directory=str(Path(entry.path).parent),
                size=stat.st_size if not is_directory else 0,
                created_time=datetime.datetime.fromtimestamp(stat.st_ctime),
                modified_time=datetime.datetime.fromtimestamp(stat.st_mtime),
                accessed_time=datetime.datetime.fromtimestamp(stat.st_atime)
            )

            if not is_directory:
                file_info.file_hash = self.calculate_file_hash(entry.path)
                try:
                    file_info.mime_type = magic.from_file(entry.path, mime=True)
                except Exception as e:
                    self.logger.warning(f"Failed to determine mime type for {entry.path}: {e}")
                try:
                    if file_info.mime_type and 'text' in file_info.mime_type:
                        with open(entry.path, 'r', encoding='utf-8') as f:
                            content = f.read()
                    else:        
                        self.extracted_text.file_path = entry.path
                        content = self.extracted_text.extract_text()
                    print(f"extracted content {content} of {file_info.path}")    
                    sensitive_data, patterns, findings = self.detect_sensitive_data(content)
                    if sensitive_data:
                        file_info.has_sensitive_data = True
                        file_info.sensitive_data = sensitive_data
                        file_info.patterns = patterns
                        self.store_sensitive_findings(entry.path, findings)
                except UnicodeDecodeError:
                        self.logger.debug(f"Unable to decode {entry.path} as UTF-8")

            self.update_database(file_info)
            return file_info

        except Exception as e:
            self.logger.error(f"Error processing {entry.path}: {e}")
            return None

    def store_sensitive_findings(self, file_path: str, findings: List[Dict]) -> None:
        with self.get_db_connection() as conn:
            cursor = conn.cursor()
            try:
                cursor.execute("SELECT id FROM files WHERE path = ?", (file_path,))
                result = cursor.fetchone()
                if result:
                    file_id = result[0]
                    cursor.executemany("""
                        INSERT INTO sensitive_data (file_id, file_path, data_type, line_number, matched_content)
                        VALUES (?, ?, ?, ?, ?)
                    """, [(file_id, file_path, f['pattern'], f['line_number'], f['matched_content']) 
                         for f in findings])
            except sqlite3.Error as e:
                self.logger.error(f"Database error while storing findings for {file_path}: {e}")

    def update_database(self, file_info: FileInfo) -> None:
        with self.get_db_connection() as conn:
            try:
                conn.execute('''
                    INSERT OR REPLACE INTO files 
                    (path, file_hash, mime_type, size, created_time, modified_time, 
                     accessed_time, is_directory, parent_directory, has_sensitive_data, 
                     patterns, sensitive_data)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                ''', (
                    file_info.path, file_info.file_hash, file_info.mime_type,
                    file_info.size, file_info.created_time, file_info.modified_time,
                    file_info.accessed_time, file_info.is_directory,
                    file_info.parent_directory, file_info.has_sensitive_data,
                    file_info.patterns, file_info.sensitive_data
                ))

                with self.lock:
                    self.stats['processed_files'] += 1
                    if not file_info.is_directory:
                        self.stats['total_size'] += file_info.size
                    if file_info.has_sensitive_data:
                        self.stats['sensitive_files'] += 1
            except sqlite3.Error as e:
                self.logger.error(f"Database error while updating file info for {file_info.path}: {e}")

    def index_filesystem(self, root_path: str) -> None:
        self.logger.info(f"Starting indexing of {root_path}")
        
        entries = list(self.collect_files(root_path))
        self.stats['total_files'] = len(entries)
        self.logger.info(f"Found {self.stats['total_files']} items to process")
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=os.cpu_count()) as executor:
            futures = [executor.submit(self.process_entry, entry) for entry in entries]
            
            for future in concurrent.futures.as_completed(futures):
                try:
                    future.result()
                except Exception as e:
                    self.logger.error(f"Task failed: {e}")
        
        self.print_summary()

    def collect_files(self, directory: str) -> List[os.DirEntry]:
        entries = []
        try:
            with os.scandir(directory) as scanner:
                for entry in scanner:
                    entries.append(entry)
                    if entry.is_dir(follow_symlinks=False):
                        entries.extend(self.collect_files(entry.path))
        except PermissionError:
            self.logger.warning(f"Permission denied: {directory}")
        except Exception as e:
            self.logger.error(f"Error scanning {directory}: {e}")
        return entries

    def print_summary(self) -> None:
        self.logger.info("\nIndexing Summary:")
        self.logger.info(f"Total files indexed: {self.stats['total_files']}")
        self.logger.info(f"Files containing sensitive data: {self.stats['sensitive_files']}")
        self.logger.info(f"Total data indexed: {self.stats['total_size'] / (1024*1024*1024):.2f} GB")

def main():
    start_time = time.time()
    
    try:
        db_path = "file_index.db"
        target_directory = r"C:\Users\Shreshth Graak\reticen\VIVEK\dlp\fingerprinting\sensitivefiles"
        
        indexer = FileSystemIndexer(db_path)
        indexer.index_filesystem(target_directory)
        
        print(f"\nTotal time taken: {time.time() - start_time:.2f} seconds")
        
    except KeyboardInterrupt:
        print("\nIndexing interrupted by user.")
        sys.exit(1)
    except Exception as e:
        print(f"\nError during indexing: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    main()
