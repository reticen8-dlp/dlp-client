import os,sys
import json
import argparse
import logging
import re
import numpy as np
from typing import Set, List, Dict, Tuple
from numpy.linalg import norm
import itertools
from contextlib import contextmanager
from multiprocessing import Pool, Manager
import pdfminer.high_level
import easyocr
import csv
import pandas as pd
from pptx import Presentation
from docx import Document
import sqlite3

import tempfile
from pdf2image import convert_from_path
import win32com.client

import magic
import difflib
import multiprocessing


logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s]: %(message)s',
    handlers=[
        logging.FileHandler('dlp_fingerprint_optimized.log'),
        logging.StreamHandler()
    ]
)



# Worker function for processing files
def process_file_worker(file_path, queue):
    try:

        extractor = TextExtractor(file_path)
        
        mime_type = magic.from_file(file_path, mime=True)
        if mime_type and 'text' in mime_type:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
        else:
            content = extractor.extract_text()
            
        if not content:
            logging.warning(f"No content extracted from {file_path}")
            return
        print(f"===========Processing file: {file_path} its content is: {content}-----------------")
        normalized = re.sub(r'[\s\W_]+', '', content.lower())
 
        
        queue.put((file_path, normalized))
    except Exception as e:
        logging.error(f"Error processing file {file_path}: {e}")


class TextExtractor:
    def __init__(self, file_path: str):
        self.file_path = file_path
        self.reader_ocr = None  # Initialize on demand to avoid pickling issues

    def _get_ocr_reader(self):
        if self.reader_ocr is None:
            self.reader_ocr = easyocr.Reader(['en'], gpu=False)
        return self.reader_ocr

    def extract_text(self) -> str:
        ext = os.path.splitext(self.file_path)[1].lower()
        try:
            if ext == '.pdf':
                return self.extract_pdf()
            elif ext == '.docx':
                return self.extract_docx()
            elif ext in ['.xls', '.xlsx']:
                return self.extract_excel()
            elif ext in ['.csv']:
                return self.extract_csv()
            elif ext in ['.pptx', '.ppt']:
                return self.extract_text_from_ppt()
            elif ext in [".jpg", ".jpeg", ".png", ".bmp", ".gif"]:
                return self.extract_image()
            elif ext in [".txt", ".md", ".html", ".htm", ".xml", ".json"]:
                return self.extract_plain_text()
            else:
                return ""
        except Exception as e:
            logging.error(f"Error processing {self.file_path}: {e}")
            return ""

    

    def extract_pdf(self) -> str:
        try:
            # Try extracting text with pdfminer first
            text = pdfminer.high_level.extract_text(self.file_path)
            if text.strip():
                return text

            # Fallback: Convert PDF pages to images for OCR
            images = convert_from_path(self.file_path, dpi=300)
            ocr_text = []
            reader = self._get_ocr_reader()
            for i, image in enumerate(images):
                temp_img_path = f"temp_page_{i}.png"
                image.save(temp_img_path, "PNG")
                ocr_result = reader.readtext(temp_img_path, detail=0)
                if ocr_result:
                    ocr_text.append("\n".join(ocr_result))
                os.remove(temp_img_path)
            return "\n".join(ocr_text)
        except Exception as e:
            logging.error(f"PDF extraction failed: {e}")
            return ""



    def extract_docx(self) -> str:
        try:
            doc = Document(self.file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            logging.error(f"DOCX extraction failed: {e}")
            return ""

    def extract_excel(self) -> str:
        try:
            xl = pd.ExcelFile(self.file_path)
            return "\n".join(" ".join(map(str, row)) for sheet in xl.sheet_names for row in xl.parse(sheet).values)
        except Exception as e:
            logging.error(f"Excel extraction failed: {e}")
            return ""

    def extract_csv(self) -> str:
        try:
            with open(self.file_path, newline='', encoding='utf-8') as csvfile:
                reader = csv.reader(csvfile)
                return "\n".join(" ".join(row) for row in reader)
        except Exception as e:
            logging.error(f"CSV extraction failed: {e}")
            return ""



    def extract_text_from_ppt(self): 
        """Extract text from PowerPoint files including images.""" 
        import os
        import tempfile
        import logging
        import pythoncom 

        ext = os.path.splitext(self.file_path)[1].lower() 
        text = "" 
        
        if ext == ".pptx":         
            prs = Presentation(self.file_path) 
            
            # Extract text from shapes 
            for slide in prs.slides: 
                for shape in slide.shapes: 
                    if hasattr(shape, "text"): 
                        text += shape.text + "\n" 
            
            # Try to extract images and perform OCR 
            try: 
                reader = self._get_ocr_reader()  # Assuming self.reader_ocr is an EasyOCR reader instance
                
                with tempfile.TemporaryDirectory() as temp_dir: 
                    img_count = 0 
                    for slide in prs.slides: 
                        for shape in slide.shapes: 
                            if shape.shape_type == 13:  # Picture shape type 
                                image = shape.image 
                                img_path = os.path.join(temp_dir, f"img_{img_count}.png") 
                            
                                with open(img_path, "wb") as f: 
                                    f.write(image.blob) 
                                
                                # Perform OCR 
                                ocr_result = reader.readtext(img_path, detail=0) 
                                if ocr_result: 
                                    text += "\n" + "\n".join(ocr_result) 
                                img_count += 1 
                                
            except ImportError: 
                logging.warning("EasyOCR not installed. Image text extraction will be skipped.") 
            except Exception as e: 
                logging.error(f"Error extracting images from PPTX: {e}") 
                    
        elif ext == ".ppt": 
            
            try: 
                # Initialize COM before creating PowerPoint instance
                pythoncom.CoInitialize()  # Add this line
                
                # Dispatch PowerPoint application
                powerpoint = win32com.client.Dispatch("PowerPoint.Application") 
                powerpoint.Visible = True  
                
                presentation = powerpoint.Presentations.Open(os.path.abspath(self.file_path), WithWindow=False) 
                
                # Extract text from slides 
                for slide in presentation.Slides: 
                    for shape in slide.Shapes: 
                        if shape.HasTextFrame: 
                            if shape.TextFrame.HasText: 
                                text += shape.TextFrame.TextRange.Text + "\n" 
                
                # Save as temporary PPTX to extract images 
                with tempfile.TemporaryDirectory() as temp_dir: 
                    temp_pptx = os.path.join(temp_dir, "temp.pptx") 
                    presentation.SaveAs(temp_pptx) 
                    presentation.Close() 
                    powerpoint.Quit() 
                    
                    # Now extract images from the PPTX version 
                    original_file_path = self.file_path  # store original file path 
                    self.file_path = temp_pptx 
                    text += self.extract_text_from_ppt() 
                    self.file_path = original_file_path  # restore original path
                    
            except Exception as e: 
                logging.error(f"Error processing PPT file: {e}") 
                if 'presentation' in locals(): 
                    presentation.Close() 
                if 'powerpoint' in locals(): 
                    powerpoint.Quit()
            finally:
                # Uninitialize COM when done
                pythoncom.CoUninitialize()  # Add this line
        else: 
            raise ValueError(f"Unsupported PowerPoint file type: {ext}") 
            
        return text


    def extract_image(self) -> str:
        try:
            import cv2
            import numpy as np
            
            # First try to load the image with OpenCV
            try:
                img = cv2.imread(self.file_path)
                if img is None:
                    logging.error(f"Unable to load image with OpenCV: {self.file_path}")
                    raise ValueError(f"OpenCV could not load image: {self.file_path}")
                    
                # Convert from BGR to RGB (EasyOCR expects RGB)
                img_rgb = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
                
                # Pass the loaded image to EasyOCR instead of the file path
                reader = self._get_ocr_reader()
                return "\n".join(reader.readtext(img_rgb, detail=0))
                
            except Exception as cv_error:
                logging.error(f"OpenCV loading failed: {cv_error}")
                
                # Try PIL as fallback
                try:
                    from PIL import Image
                    image = Image.open(self.file_path)
                    img_array = np.array(image)
                    
                    reader = self._get_ocr_reader()
                    return "\n".join(reader.readtext(img_array, detail=0))
                    
                except Exception as pil_error:
                    logging.error(f"PIL loading failed: {pil_error}")
                    raise ValueError(f"All image loading methods failed for {self.file_path}")
                    
        except Exception as e:
            logging.error(f"Image extraction failed: {e}")
            return ""

    def extract_plain_text(self) -> str:
        try:
            with open(self.file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logging.error(f"Plain text extraction failed: {e}")
            return ""

class EnhancedFingerprinter:
    def __init__(self, db_path: str = "Proprium_dlp.db"):
       self.setup_database(db_path)
       self.index_data = {}
       
    @contextmanager
    def get_db_connection(self, db_path: str = "Proprium_dlp.db"): 
        """Database connection context manager with proper initialization""" 
        conn = sqlite3.connect(db_path) 
        conn.execute('PRAGMA journal_mode=WAL') 
        conn.execute('PRAGMA synchronous=NORMAL') 
    
        try: 
            yield conn 
        except Exception as e: 
            conn.rollback() 
            raise 
        else: 
            conn.commit() 
        finally: 
            conn.close()


    def setup_database(self,db_path) -> None:
        """Initialize SQLite database with optimized schema"""
        with sqlite3.connect(db_path) as conn:  # Use direct connection for setup
            conn.executescript('''
                CREATE TABLE IF NOT EXISTS file_Fingerprint(
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    path TEXT UNIQUE,
                    normalized_text TEXT
                );
                
            ''')
            conn.commit()


    def normalize_text(self, text: str) -> str:
        return re.sub(r'[\s\W_]+', '', text.lower())
    
    def direct_match_percentage(self,str1: str, str2: str) -> float:
        """
        Checks if str2 is a subsequence of str1 and calculates the match percentage.

        :param str1: The main reference string (potential sensitive data).
        :param str2: The search string (potential leaked data).
        :return: Match percentage (0.0 to 100.0) indicating how much of str2 appears in str1 in order.
        """
        i, j = 0,0 # Pointers for str1 and str2
        while i < len(str1) and j < len(str2):
            if str1[i] == str2[j]:  # If chars match, move str2 pointer
                j += 1
               
            i += 1  # Always move str1 pointer
        
        match_percentage = (j / len(str1)) * 100 
        return match_percentage

    def compute_difflib_ratio(self, text1: str, text2: str):
        matcher = difflib.SequenceMatcher(None, text1, text2)
        diff_sim = difflib.SequenceMatcher(None, text1, text2).ratio()

        matching_blocks = matcher.get_matching_blocks()
    
        # Sum up the lengths of the matching blocks
        match_length = sum(block.size for block in matching_blocks)
        
        # Calculate the percentage of str1 that is found in str2
        match_percentage = (match_length / len(text1)) * 100 if len(text2) > 0 else 0
        return match_percentage ,diff_sim
    

    def build_index(self, folder_path: str, db_path: str) -> None:
        index_data = {}
        logging.info(f"Building index from folder: {folder_path}")
        file_paths = []
        for root, _, files in os.walk(folder_path):
            for file in files:
                file_paths.append(os.path.join(root, file))
        logging.info(f"Found files: {file_paths}")

        with Manager() as manager:
            queue = manager.Queue()
            with Pool(processes=os.cpu_count()) as pool:
                args = [(file, queue) for file in file_paths]
                pool.starmap(process_file_worker, args)
            
            # Collect results from the queue
            with self.get_db_connection(db_path) as conn:
                cursor = conn.cursor()
                while not queue.empty():
                    file_path, normalized = queue.get()
                    try:
                        cursor.execute('''INSERT OR REPLACE INTO file_Fingerprint (path, normalized_text) VALUES (?, ?)''', (file_path, normalized))
                        logging.info(f"Successfully processed  files")
                    except sqlite3.Error as e:
                        logging.error(f"Database error while updating file info for {file_path}: {e}")

           

        
   
        

    def load_index(self,db_path: str) -> None:
        index_data = {}
        with self.get_db_connection(db_path) as conn:
            cursor = conn.cursor()
            try:
                cursor.execute("SELECT path, normalized_text FROM file_Fingerprint")
                for path, normalized in cursor.fetchall():
                    index_data[path] = normalized
            except sqlite3.Error as e:
                self.logger.error(f"Database error while fetching data: {e}")
        return index_data
    
    def scan_file(self, text, index_data) -> List[Dict]:
        if not index_data:
            print("No index data loaded. Call load_index() first.")
            return []


        normalized_input = self.normalize_text(text)

        matches = []
  
        for file_path, file_data in index_data.items():
  
            file_norm = file_data
            # Direct substring check
            direct_match = normalized_input in file_norm
          
            direct_match_percentage = self.direct_match_percentage(normalized_input, file_norm)
            match_percent,diff_sim = self.compute_difflib_ratio(normalized_input, file_norm)
            matches.append({

                'file_path': file_path,
                'direct-match' : direct_match,
                'direct_match_percentage': direct_match_percentage,
                'partial_similarity': match_percent,
               
            })
        
        return sorted(matches, key=lambda x: x['partial_similarity'], reverse=True)

multiprocessing.set_start_method('spawn', force=True)
def main():
    
    parser = argparse.ArgumentParser(description="Optimized DLP Fingerprinting System")
    subparsers = parser.add_subparsers(dest="command", required=True)
    
    # Build subcommand
    parser_build = subparsers.add_parser("build", help="Build fingerprint index")
    parser_build.add_argument("--folder", required=True, help="Folder containing sensitive files to index")
    parser_build.add_argument("--db", default="Proprium_dlp.db", help="Output index file")
    
    # Scan subcommand
    parser_scan = subparsers.add_parser("scan", help="Scan a file against the index")
    parser_scan.add_argument("--text", required=True, help="text to scan")
    parser_scan.add_argument("--db", default="Proprium_dlp.db", help="Index file to use")
    
    args = parser.parse_args()
    
    
    
    if args.command == "build":
        print("scaning")
        fingerprinter = EnhancedFingerprinter(args.db)  
        print("............")  
        fingerprinter.build_index(args.folder, args.db)
    elif args.command == "scan":
        # for root, _, files in os.walk(args.folder):
        #     for file in files:
        #         input_file = os.path.join(root, file)
        fingerprinter = EnhancedFingerprinter(args.db)  
        fingerprinter.load_index(args.db)
        matches = fingerprinter.scan_file(args.text)
        # os.remove(input_file)

        
        if matches:
            print("\nMatches found:")
            for idx, match in enumerate(matches[:5], 1):  # Show top 5 matches
                print(f"{idx}. {match['file_path']}")
                print(f"   direct-match similarity: {match['direct-match']}")
                print(f"   partial similarity: {match['partial_similarity']:.2f}%")
                print(f"   direct-match percentage: {match['direct_match_percentage']}%")
            
                if match['direct-match']:
                    print("   Direct match found!")
                print()
        else:
            print("No matches found.")



if __name__ == "__main__":
    # FOR EXE:
      # if sys.platform == "win32":
      #   multiprocessing.freeze_support() 
    # Set sharing strategy for multiprocessing
    multiprocessing.set_start_method('spawn', force=True)
    
    main()




    # Tokenization
    # Anonymization:
    # Use of Differential Privacy
    # Tokenization, anonymization, and privacy-preserving comparison methods are recommended when dealing with sensitive information.
